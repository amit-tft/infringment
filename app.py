import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain_community.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
import pytesseract
from PIL import Image
from docx import Document
import pandas as pd
import textract
from pptx import Presentation
import tensorflow as tf
import numpy as np
import openai
import os

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Load the pre-trained MobileNetV2 model and the ImageNet class labels
model = tf.keras.applications.MobileNetV2(weights="imagenet")
decode_predictions = tf.keras.applications.mobilenet_v2.decode_predictions
preprocess_input = tf.keras.applications.mobilenet_v2.preprocess_input

# Set up OpenAI API key for ChatGPT


def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_image_text(image_files):
    text = ""
    for image in image_files:
        img = Image.open(image)
        text += pytesseract.image_to_string(img)
    return text

def classify_image(image_path):
    img = Image.open(image_path).resize((224, 224))
    img_array = np.array(img)
    img_array = np.expand_dims(img_array, axis=0)
    img_array = preprocess_input(img_array)
    
    predictions = model.predict(img_array)
    decoded_predictions = decode_predictions(predictions, top=3)[0]
    
    return decoded_predictions

def display_and_classify_images(image_files):
    for image in image_files:
        img = Image.open(image)
        st.image(img, caption=image.name, use_column_width=True)
        predictions = classify_image(image)
        st.write("Predictions:")
        for i, (imagenet_id, label, score) in enumerate(predictions):
            st.write(f"{i + 1}: {label} ({score:.2f})")
        explain_predictions_with_chatgpt(predictions)

def explain_predictions_with_chatgpt(predictions):
    prediction_text = "\n".join([f"{i + 1}: {label} ({score:.2f})" for i, (_, label, score) in enumerate(predictions)])
    prompt = f"Here are the top predictions from an image classification model:\n{prediction_text}\nCan you provide more details about these predictions?"

    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=150
    )
    
    explanation = response.choices[0]['message']['content'].strip()
    st.write("ChatGPT Explanation:")
    st.write(explanation)

def get_docx_text(docx_files):
    text = ""
    for doc in docx_files:
        doc_obj = Document(doc)
        for para in doc_obj.paragraphs:
            text += para.text
    return text

def get_excel_text(excel_files):
    text = ""
    for excel in excel_files:
        df = pd.read_excel(excel)
        for col in df.columns:
            text += " ".join(df[col].astype(str).tolist())
    return text

def get_csv_text(csv_files):
    text = ""
    for csv in csv_files:
        df = pd.read_csv(csv)
        for col in df.columns:
            text += " ".join(df[col].astype(str).tolist())
    return text

def get_text_file_text(text_files):
    text = ""
    for txt in text_files:
        text += txt.getvalue().decode("utf-8")
    return text

def get_ppt_text(ppt_files):
    text = ""
    for ppt in ppt_files:
        try:
            presentation = Presentation(ppt)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    elif hasattr(shape, "text_frame"):
                        text += shape.text_frame.text + "\n"
        except Exception as e:
            st.error(f"Error processing PPTX file {ppt.name}: {e}")
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    st.set_page_config(page_title="Infringment Detection Tool",
                       page_icon=":male-astronaut:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Infringment Detection Tool ")
    user_question = st.text_input("Ask a question:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your image file")
        uploaded_files = st.file_uploader(
            "Upload your files here and click on 'Process' & ask question",
            accept_multiple_files=True,
            type=["pdf","png", "jpg", "jpeg","docx"]
        )
        if st.button("Process"):
            with st.spinner("Processing"):
                pdf_docs = [f for f in uploaded_files if f.type == "application/pdf"]
                image_files = [f for f in uploaded_files if f.type in ["image/png", "image/jpeg"]]
                docx_files = [f for f in uploaded_files if f.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]
                csv_files = [f for f in uploaded_files if f.type == "text/csv"]
                excel_files = [f for f in uploaded_files if f.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]
                text_files = [f for f in uploaded_files if f.type == "text/plain"]
                ppt_files = [f for f in uploaded_files if f.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"]

                raw_text = ""
                if pdf_docs:
                    raw_text += get_pdf_text(pdf_docs)
                if image_files:
                    display_and_classify_images(image_files)
                if docx_files:
                    raw_text += get_docx_text(docx_files)
                if excel_files:
                    raw_text += get_excel_text(excel_files)
                if csv_files:
                    raw_text += get_csv_text(csv_files)
                if text_files:
                    raw_text += get_text_file_text(text_files)
                if ppt_files:
                    raw_text += get_ppt_text(ppt_files)

                if not raw_text.strip() and not image_files:
                    st.error("No text or images could be extracted from the uploaded files.")
                    return

                if raw_text.strip():
                    text_chunks = get_text_chunks(raw_text)
                    if not text_chunks:
                        st.error("No text chunks could be created. Please check the uploaded files.")
                        return

                    vectorstore = get_vectorstore(text_chunks)
                    st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()
